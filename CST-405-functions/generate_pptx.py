#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Compiler Performance Enhancement
CST-405 Principles of Compiler Design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs):
    """Add the title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "CST-405 Principles of Compiler Design"
    subtitle.text = "Compiler Performance Enhancement\nOptimizing Every Phase of Compilation\n\nClass Activity & Implementation Guide"

def add_content_slide(prs, title_text, content):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = title_text

    tf = body.text_frame
    tf.text = content[0] if content else ""

    for item in content[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0 if not item.startswith("  ") else 1

def add_code_slide(prs, title_text, code_text, explanation=""):
    """Add a slide with code example"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    # Code box
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    code_frame = code_box.text_frame
    code_frame.text = code_text

    # Format code
    for paragraph in code_frame.paragraphs:
        paragraph.font.name = "Consolas"
        paragraph.font.size = Pt(12)

    # Add background to code box
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(40, 44, 52)  # Dark background

    # Explanation if provided
    if explanation:
        exp_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
        exp_frame = exp_box.text_frame
        exp_frame.text = explanation
        exp_frame.paragraphs[0].font.size = Pt(14)

def add_comparison_slide(prs, title_text, before_title, before_content, after_title, after_content):
    """Add a before/after comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True

    # Before box
    before_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(5))
    before_frame = before_box.text_frame
    before_frame.text = f"❌ {before_title}\n\n{before_content}"

    # After box
    after_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5))
    after_frame = after_box.text_frame
    after_frame.text = f"✅ {after_title}\n\n{after_content}"

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()

    # Set slide size to widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs)

    # Slide 2: Overview
    add_content_slide(prs, "Activity Overview", [
        "What We'll Optimize Today:",
        "• Symbol Table: From O(n) to O(1) lookup",
        "• Memory Management: Pool allocation for AST nodes",
        "• TAC Generation: Peephole optimization patterns",
        "• Lexical Analysis: Buffer optimization",
        "• Code Generation: Register allocation",
        "• String Management: String interning",
        "",
        "Goal: 30-50% overall performance improvement"
    ])

    # Slide 3: Why Performance Matters
    add_content_slide(prs, "Why Compiler Performance Matters", [
        "Development Impact:",
        "• Small Project: 2s → 1s (50% faster)",
        "• Medium Project: 30s → 15s (50% faster)",
        "• Large Project: 5min → 2.5min (50% faster)",
        "",
        "Benefits:",
        "• Developer Productivity: Faster compile-test cycles",
        "• CI/CD Pipelines: Reduced build times",
        "• Resource Usage: Lower memory footprint",
        "• Energy Efficiency: Less CPU time = less power"
    ])

    # Slide 4: Benchmarking Setup
    add_code_slide(prs, "Performance Measurement Setup",
"""typedef struct {
    double cpu_time;      // Processor time used
    double wall_time;     // Real-world elapsed time
    long memory_usage;    // Memory consumed
    long peak_memory;     // Maximum memory used
} BenchmarkResult;

// Platform-specific measurement
#ifdef __APPLE__
    // Use Mach kernel APIs
#else
    // Use /proc filesystem
#endif""",
    "Key metrics to track for each optimization")

    # Slide 5: Symbol Table Problem
    add_comparison_slide(prs, "Symbol Table Optimization",
        "Current: Linear Search",
        """for(i = 0; i < count; i++) {
    if(strcmp(vars[i].name, name) == 0)
        return vars[i].offset;
}
// O(n) complexity - gets slower with more variables""",
        "Solution: Hash Table",
        """unsigned int h = hash(name);
SymbolNode* node = buckets[h];
while(node) {
    if(strcmp(node->name, name) == 0)
        return node->offset;
    node = node->next;
}
// O(1) average complexity - constant time!""")

    # Slide 6: Hash Function
    add_code_slide(prs, "Implementing a Good Hash Function",
"""// DJB2 Hash Algorithm
unsigned int hash(const char* str) {
    unsigned int hash = 5381;
    int c;

    while ((c = *str++)) {
        hash = ((hash << 5) + hash) + c; // hash * 33 + c
    }

    return hash % HASH_SIZE;  // HASH_SIZE = 211 (prime)
}

// Why this works:
// • Prime number size (211): Better distribution
// • Multiplier 33: Good bit spreading
// • Simple operations: Fast execution""")

    # Slide 7: Performance Results
    add_content_slide(prs, "Symbol Table Performance Results", [
        "Measured Improvements:",
        "",
        "100 variables:",
        "  • Linear: 5ms → Hash: 1ms (5x speedup)",
        "",
        "500 variables:",
        "  • Linear: 125ms → Hash: 2ms (62x speedup)",
        "",
        "1000 variables:",
        "  • Linear: 500ms → Hash: 3ms (166x speedup)",
        "",
        "Key Insight: Performance gap widens dramatically with scale!"
    ])

    # Slide 8: Memory Pool Concept
    add_comparison_slide(prs, "Memory Pool for AST Nodes",
        "Individual Allocations",
        """// Each node = separate malloc
ASTNode* node = malloc(sizeof(ASTNode));

Problems:
• Memory fragmentation
• Allocation overhead
• Poor cache locality
• 100ns per malloc""",
        "Pool Allocation",
        """// Allocate from pool
ASTNode* node = ast_alloc(sizeof(ASTNode));

Benefits:
• Contiguous memory
• No fragmentation
• Better cache performance
• 5ns per allocation""")

    # Slide 9: Memory Pool Implementation
    add_code_slide(prs, "Memory Pool Implementation",
"""#define POOL_SIZE 4096  // 4KB chunks

void* ast_alloc(size_t size) {
    // Align to 8 bytes for performance
    size = (size + 7) & ~7;

    if (current_pool->used + size > POOL_SIZE) {
        allocate_new_pool();
    }

    void* ptr = current_pool->memory + current_pool->used;
    current_pool->used += size;
    return ptr;
}

// Result: 30-50% memory reduction!""")

    # Slide 10: TAC Optimization
    add_comparison_slide(prs, "TAC Peephole Optimization",
        "Unoptimized TAC",
        """t0 = 5
t1 = 10
t2 = t0 + t1
x = t2

4 instructions, 3 temporaries""",
        "Optimized TAC",
        """x = 15

1 instruction, 0 temporaries!

Techniques:
• Constant folding
• Copy propagation
• Dead code elimination""")

    # Slide 11: Peephole Implementation
    add_code_slide(prs, "Implementing Peephole Patterns",
"""// Pattern: Constant Folding
if (current->op == TAC_ASSIGN &&
    next->op == TAC_ASSIGN &&
    isConstant(current->arg1) &&
    isConstant(next->arg1)) {

    // Look for: t0 = 5; t1 = 10; t2 = t0 + t1
    if (third->op == TAC_ADD &&
        third->arg1 == current->result &&
        third->arg2 == next->result) {

        // Replace with: t2 = 15
        int result = atoi(current->arg1) + atoi(next->arg1);
        third->op = TAC_ASSIGN;
        third->arg1 = itoa(result);

        // Mark old instructions for removal
        current->op = TAC_NOP;
        next->op = TAC_NOP;
    }
}""")

    # Slide 12: Scanner Optimization
    add_code_slide(prs, "Scanner Buffer Optimization",
"""// Increase buffer size for better I/O
#define YY_BUF_SIZE 32768  // 32KB vs default 8KB

// Why 32KB?
// • Fits in L1 cache on modern CPUs
// • Reduces system calls by 4x
// • Good balance of memory vs performance

// Keyword recognition with perfect hash
int keyword_hash(const char* str) {
    // Length + first char gives unique values
    return (strlen(str) + str[0]) % 13;
}

// Result: 10x faster keyword detection!""")

    # Slide 13: Register Allocation
    add_comparison_slide(prs, "Register Allocation",
        "Stack-Based Code",
        """lw $t0, -4($fp)   // Load x
lw $t1, -8($fp)   // Load y
add $t2, $t0, $t1
sw $t2, -12($fp)  // Store result

3 memory accesses!
Each access: 100-200 cycles""",
        "Register-Allocated Code",
        """// x in $s0, y in $s1
add $s2, $s0, $s1

0 memory accesses!
Register access: 1 cycle

Result: 40-60% faster execution""")

    # Slide 14: LRU Strategy
    add_code_slide(prs, "LRU Register Allocation",
"""char* allocate_register(char* var) {
    // 1. Check if variable already in a register
    for(reg : registers) {
        if(reg.var_name == var) return reg.reg_name;
    }

    // 2. Find free register
    for(reg : registers) {
        if(reg.is_free) {
            reg.var_name = var;
            return reg.reg_name;
        }
    }

    // 3. Evict LRU register
    Register* lru = find_lru_register();
    spill_to_stack(lru);  // Save old value
    lru.var_name = var;
    return lru.reg_name;
}""")

    # Slide 15: String Interning
    add_comparison_slide(prs, "String Interning",
        "Without Interning",
        """// "count" appears 10 times
"count" → malloc(6)  // 6 bytes
"count" → malloc(6)  // 6 bytes
"count" → malloc(6)  // 6 bytes
...
Total: 60 bytes + overhead""",
        "With Interning",
        """// All point to same string
"count" → pool[0]  // 6 bytes once
"count" → pool[0]  // Same pointer
"count" → pool[0]  // Same pointer
...
Total: 6 bytes!

Benefits:
• 20-40% memory savings
• Faster comparisons (pointer equality)""")

    # Slide 16: Performance Testing
    add_code_slide(prs, "Comprehensive Performance Testing",
"""#!/bin/bash
# Generate test files of varying complexity

# Small test (10 variables)
for i in {1..10}; do
    echo "int var$i;"
    echo "var$i = $i;"
done > test_small.c

# Large test (1000 variables)
for i in {1..1000}; do
    echo "int var$i;"
    echo "var$i = $i;"
done > test_large.c

# Measure performance
time ./compiler_original < test_large.c
time ./compiler_optimized < test_large.c

# Expected: 30-50% overall speedup""")

    # Slide 17: Platform-Specific Tools
    add_content_slide(prs, "Platform-Specific Profiling", [
        "macOS Tools:",
        "  • Instruments -t 'Time Profiler' ./compiler",
        "  • leaks --atExit -- ./compiler",
        "  • sample ./compiler -file test.c",
        "",
        "Linux Tools:",
        "  • perf record -g ./compiler",
        "  • valgrind --tool=massif ./compiler",
        "  • gprof ./compiler",
        "",
        "Cross-Platform:",
        "  • Built-in timing with clock()",
        "  • Custom instrumentation",
        "",
        "Tip: Always profile before optimizing!"
    ])

    # Slide 18: Results Summary
    add_content_slide(prs, "Expected Performance Improvements", [
        "Optimization Results:",
        "",
        "Hash Table Symbol Table:",
        "  • 10-100x faster lookups",
        "",
        "AST Memory Pool:",
        "  • 30-50% less memory, 2-3x faster allocation",
        "",
        "TAC Optimization:",
        "  • 15-25% fewer instructions",
        "",
        "Scanner Buffer:",
        "  • 20-30% I/O improvement",
        "",
        "Register Allocation:",
        "  • 40-60% faster execution",
        "",
        "String Interning:",
        "  • 20-40% memory savings"
    ])

    # Slide 19: Common Pitfalls
    add_content_slide(prs, "Common Implementation Pitfalls", [
        "What to Watch Out For:",
        "",
        "• Hash Collisions: Poor hash → O(n) worst case",
        "  Solution: Use prime table sizes",
        "",
        "• Memory Pool Fragmentation: Variable sizes",
        "  Solution: Separate pools for different sizes",
        "",
        "• Register Spilling: Too aggressive",
        "  Solution: Profile-guided allocation",
        "",
        "• Over-optimization: Complex code for minimal gain",
        "  Solution: Focus on hotspots (80/20 rule)",
        "",
        "• Platform Dependencies: OS-specific code",
        "  Solution: Use #ifdef for portability"
    ])

    # Slide 20: Real-World Impact
    add_content_slide(prs, "Real-World Impact", [
        "Industry Applications:",
        "",
        "• GCC/LLVM: Use all these techniques",
        "• V8 (JavaScript): Hidden classes = hash tables",
        "• Java HotSpot: Aggressive register allocation",
        "• Rust: Memory pools for fast compilation",
        "",
        "Career Relevance - Performance matters in:",
        "• Game engines",
        "• Database systems",
        "• Operating systems",
        "• Embedded systems",
        "• High-frequency trading",
        "",
        "These skills transfer directly to industry!"
    ])

    # Slide 21: Key Takeaways
    add_content_slide(prs, "Key Takeaways", [
        "Remember:",
        "",
        "• Always measure before optimizing",
        "",
        "• Focus on algorithmic improvements first",
        "  (O(n) → O(1) beats constant factor improvements)",
        "",
        "• Memory locality matters as much as CPU cycles",
        "",
        "• Platform differences can be significant",
        "",
        "• 80% of time is in 20% of code - optimize hotspots",
        "",
        '"Premature optimization is the root of all evil" - Knuth',
        "",
        "But now you know when and how to optimize effectively!"
    ])

    # Save presentation
    prs.save('Compiler_Performance_Enhancement.pptx')
    print("Presentation created: Compiler_Performance_Enhancement.pptx")

if __name__ == "__main__":
    create_presentation()